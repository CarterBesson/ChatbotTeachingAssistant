#ChatGPT was used to create sections of this code
import logging
from io import BytesIO
from pathlib import Path

import magic  # Requires the 'python-magic' library

# Import libraries for extracting text from various file types
import pypdf
import docx
from bs4 import BeautifulSoup
from pptx import Presentation
from striprtf.striprtf import rtf_to_text

# Import LangChain's text splitter
from langchain_text_splitters import RecursiveCharacterTextSplitter, TokenTextSplitter

# Configure logging
logging.basicConfig(level=logging.INFO)


def detect_mime_type(file_bytes):
    """
    Detects the MIME type of a file using magic numbers.
    """
    try:
        mime = magic.from_buffer(file_bytes, mime=True)
        return mime
    except Exception as e:
        logging.error(f"Error detecting MIME type: {e}")
        return None


def clean_text(text):
    """
    Cleans the text by removing unwanted characters or formatting.
    """
    # Remove zero-width spaces
    return text.replace('\u200b', '')


def chunk_text(text, method='recursive', chunk_size=1000, chunk_overlap=200):
    """
    Splits text into chunks using LangChain's text splitters.

    Parameters:
    - method: 'recursive' or 'token'
    - chunk_size: Size of each chunk (characters or tokens)
    - chunk_overlap: Overlap between chunks (characters or tokens)
    """
    text = clean_text(text)
    if method == 'token':
        text_splitter = TokenTextSplitter(
            encoding_name="cl100k_base",
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
    else:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
    return text_splitter.split_text(text)


def handle_extraction_errors(func):
    """
    Decorator to handle exceptions and logging in extraction functions.
    """
    def wrapper(file_bytes):
        try:
            return func(file_bytes)
        except Exception as e:
            logging.error(f"Error in {func.__name__}: {e}")
            return ''
    return wrapper


@handle_extraction_errors
def extract_text_from_pdf(file_bytes):
    """
    Extracts text from a PDF file.
    """
    pdf_reader = pypdf.PdfReader(BytesIO(file_bytes))
    text = ''
    for page in pdf_reader.pages:
        text += page.extract_text() or ''
    return text


@handle_extraction_errors
def extract_text_from_docx(file_bytes):
    """
    Extracts text from a DOCX file.
    """
    document = docx.Document(BytesIO(file_bytes))
    return '\n'.join([para.text for para in document.paragraphs])


@handle_extraction_errors
def extract_text_from_txt(file_bytes):
    """
    Extracts text from a TXT file.
    """
    return file_bytes.decode('utf-8', errors='ignore')


@handle_extraction_errors
def extract_text_from_pptx(file_bytes):
    """
    Extracts text from a PPTX file.
    """
    presentation = Presentation(BytesIO(file_bytes))
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return '\n'.join(text_runs)


@handle_extraction_errors
def extract_text_from_html(file_bytes):
    """
    Extracts text from an HTML file.
    """
    soup = BeautifulSoup(file_bytes, 'html.parser')
    return soup.get_text(separator='\n')


@handle_extraction_errors
def extract_text_from_rtf(file_bytes):
    """
    Extracts text from an RTF file.
    """
    return rtf_to_text(file_bytes.decode('utf-8', errors='ignore'))


# Mappings from extensions and MIME types to extractor functions
extension_to_extractor = {
    '.pdf': extract_text_from_pdf,
    '.docx': extract_text_from_docx,
    '.txt': extract_text_from_txt,
    '.pptx': extract_text_from_pptx,
    '.html': extract_text_from_html,
    '.htm': extract_text_from_html,
    '.rtf': extract_text_from_rtf,
    # Add more mappings as needed
}

mime_type_to_extractor = {
    'application/pdf': extract_text_from_pdf,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': extract_text_from_docx,
    'text/plain': extract_text_from_txt,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': extract_text_from_pptx,
    'text/html': extract_text_from_html,
    'application/rtf': extract_text_from_rtf,
    # Add more mappings as needed
}


def process_file(file_bytes, filename):
    """
    Main function to process a file: detects file type, extracts text, and chunks it.
    """
    extension = Path(filename).suffix.lower()
    extractor = extension_to_extractor.get(extension)

    if extractor:
        logging.info(f"Using extractor for file extension: '{extension}'")
    else:
        detected_mime = detect_mime_type(file_bytes)
        extractor = mime_type_to_extractor.get(detected_mime)
        if extractor:
            logging.info(f"Using extractor for MIME type: '{detected_mime}'")
        else:
            error_message = (f"No extractor found for file '{filename}' "
                             f"with extension '{extension}' and MIME type '{detected_mime}'.")
            logging.error(error_message)
            return None

    text = extractor(file_bytes)
    if not text:
        error_message = f"No text extracted from file '{filename}'."
        logging.error(error_message)
        return None

    # Use LangChain's text splitter
    chunks = chunk_text(text, method='token', chunk_size=1000, chunk_overlap=200)
    return chunks